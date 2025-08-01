#!/usr/bin/env python3
"""
PowerPoint to Word Document Generator using Claude API
Compatible with both WSL2 and Windows Command Prompt
"""

import os
import sys
import json
import asyncio
import logging
from pathlib import Path
from typing import Dict, List, Optional
from datetime import datetime

# Third-party imports
import aiohttp
from pptx import Presentation
from docxtpl import DocxTemplate
from flask import Flask, request, jsonify, send_file
from werkzeug.utils import secure_filename

# Configure logging
logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')
logger = logging.getLogger(__name__)


# Configuration
class Config:
    CLAUDE_API_KEY = os.environ.get('CLAUDE_API_KEY', 'your-api-key-here')
    CLAUDE_API_URL = "https://api.anthropic.com/v1/messages"
    CLAUDE_MODEL = "claude-3-sonnet-20240229"  # Updated model name

    UPLOAD_FOLDER = Path('uploads')
    OUTPUT_FOLDER = Path('outputs')
    TEMPLATE_FOLDER = Path('templates')

    MAX_TOKENS = 4000
    ALLOWED_EXTENSIONS = {'pptx'}

    # Document metadata
    DEFAULT_USER_NAME = os.environ.get('USER_NAME', os.environ.get('USERNAME', 'System User'))
    MANGO_ID_PREFIX = os.environ.get('MANGO_ID_PREFIX', 'MANGO-')

    # Create directories if they don't exist
    for folder in [UPLOAD_FOLDER, OUTPUT_FOLDER, TEMPLATE_FOLDER]:
        folder.mkdir(exist_ok=True)


class SOPDocumentConfig:
    """Configuration for SOP document generation"""

    # Default Claude prompt for work instruction to SOP conversion
    DEFAULT_SOP_PROMPT = """# PowerPoint Work Instruction to SOP Conversion Prompt

You are an AI agent specialized in converting PowerPoint work instructions into Standard Operating Procedure (SOP) documents. Your task is to extract content from PowerPoint slides and format it according to the LS SOP template structure.

## Critical Instructions
- Extract ALL content from the PowerPoint being converted
- The PROCEDURE section should simply state: "See Mango File {{mango_pptx_id}} -- {{pptx_title}}."
- DO NOT generate or infer the Mango ID - this will be provided externally
- Focus on extracting the OBJECTIVE, RESPONSIBILITIES, and DEFINITIONS from the PowerPoint content

## Required Output Format

You MUST respond with a valid JSON object containing these exact fields:

```json
{
    "title": "[Extract the main title/subject from the PowerPoint]",
    "objective": "[Extract or infer the main objective/purpose of this procedure from the PowerPoint content]",
    "scope": "This SOP applies to all Cassette Manufacturing operations at the US Stoneham facility",
    "responsibilities": "[Extract specific roles mentioned: Technician, Senior Technician, Engineer, etc. - BE SPECIFIC to what's in the PowerPoint]",
    "definitions": "[Extract any defined terms, abbreviations, or acronyms WITH their definitions (e.g., 'DFF: Direct Filter Flow'). If none found, use 'N/A']",
    "mango_id": "",
    "pptx_title": "[Use the PowerPoint filename or main title]"
}
```

## Content Extraction Guidelines

### OBJECTIVE
- Look for slides mentioning purpose, goals, or objectives
- If not explicitly stated, infer from the overall content what this procedure accomplishes
- Keep it concise (2-3 sentences maximum)
- Focus on WHAT is being done and WHY

### RESPONSIBILITIES  
- Extract ACTUAL roles mentioned in the PowerPoint (not generic roles)
- Common roles include: Technician, Senior Technician, Engineer, Quality Assurance
- Include what each role does if specified
- If multiple roles are mentioned, list them all
- Format as: "Role: Specific responsibilities" if details are provided

### DEFINITIONS
- Look for slides with terminology, abbreviations, or glossary sections
- Extract ONLY terms that are explicitly defined in the PowerPoint
- Format as: "TERM: Definition" or "ABBREVIATION: Full meaning"
- Include multiple definitions separated by semicolons
- If no definitions are found, use "N/A"

### Important Notes
- The SCOPE is always hard-coded as shown above - do not modify
- The PROCEDURE section will always reference the Mango file - do not extract procedure steps
- Extract the document title from the PowerPoint title slide or filename
- All content must come from the PowerPoint - do not add generic information

Remember: Your ENTIRE response must be a single valid JSON object with no additional text or formatting."""


# Flask app for web interface
app = Flask(__name__)
app.config['MAX_CONTENT_LENGTH'] = 256 * 1024 * 1024  # 16MB max file size


class PowerPointProcessor:
    """Extract content from PowerPoint files"""

    @staticmethod
    def extract_text_from_pptx(file_path: Path) -> Dict[str, any]:
        """Extract all text content from a PowerPoint file"""
        try:
            prs = Presentation(file_path)
            extracted_data = {
                'filename': file_path.name,
                'slide_count': len(prs.slides),
                'slides': []
            }

            for idx, slide in enumerate(prs.slides, 1):
                slide_content = {
                    'slide_number': idx,
                    'title': '',
                    'content': []
                }

                # Extract title
                if slide.shapes.title:
                    slide_content['title'] = slide.shapes.title.text

                # Extract all text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        if shape != slide.shapes.title:
                            slide_content['content'].append(shape.text)

                extracted_data['slides'].append(slide_content)

            return extracted_data

        except Exception as e:
            logger.error(f"Error extracting from {file_path}: {str(e)}")
            raise


class ClaudeAPIClient:
    """Handle Claude API interactions"""

    def __init__(self, api_key: str):
        self.api_key = api_key
        self.headers = {
            "Content-Type": "application/json",
            "x-api-key": api_key,
            "anthropic-version": "2023-06-01"
        }

    async def process_presentation(self, pptx_content: Dict, custom_prompt: str) -> Dict:
        """Send PowerPoint content to Claude and get response"""

        # Format the extracted content for Claude
        formatted_content = self._format_pptx_content(pptx_content)

        # Combine with user's custom prompt
        full_prompt = f"{custom_prompt}\n\nPowerPoint Content:\n{formatted_content}"

        payload = {
            "model": Config.CLAUDE_MODEL,
            "messages": [
                {
                    "role": "user",
                    "content": full_prompt
                }
            ],
            "max_tokens": Config.MAX_TOKENS
        }

        async with aiohttp.ClientSession() as session:
            try:
                async with session.post(
                        Config.CLAUDE_API_URL,
                        headers=self.headers,
                        json=payload
                ) as response:
                    response_text = await response.text()

                    if response.status == 200:
                        try:
                            data = json.loads(response_text)
                            return {
                                'success': True,
                                'content': data['content'][0]['text'],
                                'usage': data.get('usage', {})
                            }
                        except (json.JSONDecodeError, KeyError) as e:
                            logger.error(f"Failed to parse API response: {e}")
                            logger.error(f"Response text: {response_text[:500]}")
                            return {
                                'success': False,
                                'error': 'Invalid API response format',
                                'details': str(e)
                            }
                    else:
                        # Log the actual error response
                        logger.error(f"Claude API error: {response.status}")
                        logger.error(f"Response: {response_text[:500]}")

                        # Try to parse error message
                        error_msg = f"API error: {response.status}"
                        try:
                            error_data = json.loads(response_text)
                            if 'error' in error_data:
                                error_msg = error_data['error'].get('message', error_msg)
                        except:
                            # If response is HTML, it's likely an auth error
                            if response_text.startswith('<'):
                                error_msg = "Authentication failed. Please check your API key."

                        return {
                            'success': False,
                            'error': error_msg,
                            'details': response_text[:500]
                        }

            except aiohttp.ClientError as e:
                logger.error(f"Network error calling Claude API: {str(e)}")
                return {
                    'success': False,
                    'error': f"Network error: {str(e)}"
                }
            except Exception as e:
                logger.error(f"Unexpected error calling Claude API: {str(e)}")
                return {
                    'success': False,
                    'error': f"Unexpected error: {str(e)}"
                }

    def _format_pptx_content(self, pptx_content: Dict) -> str:
        """Format extracted PowerPoint content for Claude"""
        formatted = f"Presentation: {pptx_content['filename']}\n"
        formatted += f"Total Slides: {pptx_content['slide_count']}\n\n"

        for slide in pptx_content['slides']:
            formatted += f"--- Slide {slide['slide_number']} ---\n"
            if slide['title']:
                formatted += f"Title: {slide['title']}\n"
            if slide['content']:
                formatted += "Content:\n"
                for content in slide['content']:
                    formatted += f"  • {content}\n"
            formatted += "\n"

        return formatted


class WordDocumentGenerator:
    """Generate Word documents from templates"""

    def __init__(self, template_path: Path):
        self.template_path = template_path
        if not self.template_path.exists():
            raise FileNotFoundError(f"Template not found: {template_path}")

    def generate_document(self, context: Dict, output_path: Path) -> Path:
        """Generate a Word document from template with given context"""
        try:
            doc = DocxTemplate(self.template_path)

            # Add timestamp and other metadata
            context['generated_date'] = datetime.now().strftime('%Y-%m-%d %H:%M:%S')
            context['generator_version'] = '1.0'

            # Render the document with context
            doc.render(context)

            # Save the document
            doc.save(output_path)
            logger.info(f"Generated document: {output_path}")

            return output_path

        except Exception as e:
            logger.error(f"Error generating document: {str(e)}")
            raise


class DocumentProcessor:
    """Main processor coordinating the entire workflow"""

    def __init__(self, api_key: str, template_path: Path, custom_prompt: str):
        self.pptx_processor = PowerPointProcessor()
        self.claude_client = ClaudeAPIClient(api_key)
        self.doc_generator = WordDocumentGenerator(template_path)
        self.custom_prompt = custom_prompt
        self.user_name = Config.DEFAULT_USER_NAME
        self.mango_prefix = ''

    async def process_single_file(self, pptx_path: Path) -> Path:
        """Process a single PowerPoint file to Word document"""

        # Step 1: Extract content from PowerPoint
        logger.info(f"Extracting content from: {pptx_path}")
        pptx_content = self.pptx_processor.extract_text_from_pptx(pptx_path)

        # Step 2: Send to Claude API
        logger.info("Sending to Claude API...")
        claude_response = await self.claude_client.process_presentation(
            pptx_content,
            self.custom_prompt
        )

        if not claude_response['success']:
            raise Exception(f"Claude API error: {claude_response['error']}")

        # Step 3: Parse Claude's response
        context = self._parse_claude_response(claude_response['content'])

        # Add source file information and user context
        context['source_file'] = pptx_path.name
        context['slide_count'] = pptx_content['slide_count']
        context['user_name'] = self.user_name

        # Generate Mango ID if prefix is provided and no ID was found
        if self.mango_prefix and not context.get('mango_pptx_id'):
            # Generate ID based on timestamp and filename
            timestamp = datetime.now().strftime('%Y%m%d%H%M%S')
            file_hash = str(hash(pptx_path.name))[-4:]
            context['mango_pptx_id'] = f"{self.mango_prefix}{timestamp}-{file_hash}"

        # If no title was extracted, use filename
        if not context.get('pptx_title'):
            context['pptx_title'] = pptx_path.stem.replace('_', ' ').title()

        # Step 4: Generate Word document
        output_filename = f"{pptx_path.stem}_SOP_{datetime.now().strftime('%Y%m%d')}.docx"
        output_path = Config.OUTPUT_FOLDER / output_filename

        return self.doc_generator.generate_document(context, output_path)

    def _parse_claude_response(self, claude_text: str) -> Dict:
        """
        Parse Claude's response into template context for LS SOP template
        """
        # Initialize context with all required template variables
        context = {
            # Template placeholders from your document
            'mango_pptx_id': '',
            'pptx_title': '',
            'date': datetime.now().strftime('%m/%d/%Y'),  # US date format
            'user_name': self.user_name or os.environ.get('USER_NAME', 'System User'),

            # SOP Sections - defaults
            'objective': '',
            'scope': 'This SOP applies to all Cassette Manufacturing operations at the US Stoneham facility',
            'responsibilities': '',
            'definitions': 'N/A',

            # The procedure section is fixed per your requirements
            'procedure_text': 'See Mango File {{mango_pptx_id}} -- {{pptx_title}}.',

            # Revision tracking - for the table
            'revisions': [
                {
                    'date': datetime.now().strftime('%m/%d/%Y'),
                    'version': '1.0',
                    'nature_of_changes': 'Initial Release',
                    'changed_by': self.user_name or os.environ.get('USER_NAME', 'System User')
                }
            ],

            # Keep full response for debugging
            'full_response': claude_text
        }

        # Parse Claude's JSON response
        try:
            # Clean the response - remove markdown code blocks if present
            cleaned_response = claude_text.strip()
            if cleaned_response.startswith('```json'):
                cleaned_response = cleaned_response[7:]  # Remove ```json
            if cleaned_response.startswith('```'):
                cleaned_response = cleaned_response[3:]  # Remove ```
            if cleaned_response.endswith('```'):
                cleaned_response = cleaned_response[:-3]  # Remove trailing ```

            parsed = json.loads(cleaned_response.strip())

            # Map Claude's response to template fields
            context['objective'] = parsed.get('objective', '')
            context['scope'] = parsed.get('scope', context['scope'])  # Use default if not provided
            context['responsibilities'] = parsed.get('responsibilities', '')
            context['definitions'] = parsed.get('definitions', 'N/A')
            context['pptx_title'] = parsed.get('pptx_title', parsed.get('title', ''))

            # Mango ID might come from Claude or be set externally
            if parsed.get('mango_id'):
                context['mango_pptx_id'] = parsed.get('mango_id')

        except json.JSONDecodeError as e:
            logger.error(f"Failed to parse Claude's JSON response: {e}")
            logger.error(f"Response was: {claude_text}")

            # Fallback: try to extract some basic info from text
            lines = claude_text.strip().split('\n')
            if lines:
                context['pptx_title'] = lines[0].strip()
                context['objective'] = "Failed to parse response. Please review manually."

        return context

    async def process_batch(self, pptx_files: List[Path]) -> List[Path]:
        """Process multiple PowerPoint files"""
        results = []
        for pptx_file in pptx_files:
            try:
                output_path = await self.process_single_file(pptx_file)
                results.append(output_path)
            except Exception as e:
                logger.error(f"Failed to process {pptx_file}: {str(e)}")
                results.append(None)

        return results


# Flask Web Interface Routes
@app.route('/')
def index():
    """SOP Generator HTML interface"""
    return '''
    <!DOCTYPE html>
    <html>
    <head>
        <title>SOP Generator - PowerPoint to Word</title>
        <style>
            body { 
                font-family: Verdana, Arial, sans-serif; 
                margin: 40px;
                background-color: #f5f5f5;
            }
            .container {
                background-color: white;
                padding: 30px;
                border-radius: 8px;
                box-shadow: 0 2px 4px rgba(0,0,0,0.1);
                max-width: 800px;
                margin: 0 auto;
            }
            h1 { color: #333; }
            .upload-area { 
                border: 2px dashed #4CAF50; 
                padding: 30px; 
                text-align: center;
                margin: 20px 0;
                background-color: #f9f9f9;
                border-radius: 4px;
            }
            .form-group {
                margin: 20px 0;
            }
            label {
                display: block;
                font-weight: bold;
                margin-bottom: 5px;
                color: #555;
            }
            input[type="text"], textarea {
                width: 100%;
                padding: 8px;
                border: 1px solid #ddd;
                border-radius: 4px;
                box-sizing: border-box;
            }
            textarea { 
                font-family: monospace; 
                font-size: 12px;
            }
            button { 
                background-color: #4CAF50; 
                color: white; 
                padding: 12px 24px;
                border: none;
                cursor: pointer;
                margin: 5px;
                border-radius: 4px;
                font-size: 16px;
            }
            button:hover { background-color: #45a049; }
            .status { 
                margin-top: 20px;
                padding: 10px;
                border-radius: 4px;
            }
            .info-box {
                background-color: #e3f2fd;
                border-left: 4px solid #2196F3;
                padding: 15px;
                margin: 20px 0;
            }
            .success { 
                background-color: #c8e6c9;
                color: #2e7d32;
            }
            .error {
                background-color: #ffcccc;
                color: #cc0000;
            }
        </style>
    </head>
    <body>
        <div class="container">
            <h1>Standard Operating Procedure Generator</h1>
            <p>Convert PowerPoint presentations to SOP Word documents using Claude AI</p>

            <div class="info-box">
                <strong>Note:</strong> The generated documents will follow the LS SOP template format.
                Document ID and name will be populated when uploaded to the document management system.
            </div>

            <form method="POST" enctype="multipart/form-data" action="/upload" id="uploadForm">
                <div class="upload-area">
                    <p><strong>Select PowerPoint files to convert:</strong></p>
                    <input type="file" name="files" multiple accept=".pptx" required>
                    <p style="color: #666; font-size: 14px;">Supports multiple .pptx files</p>
                </div>

                <div class="form-group">
                    <label for="user_name">Your Name (for revision tracking):</label>
                    <input type="text" name="user_name" id="user_name" 
                           value="''' + Config.DEFAULT_USER_NAME + '''" required>
                </div>

                <div class="form-group">
                    <label for="mango_prefix">Mango ID Prefix (optional):</label>
                    <input type="text" name="mango_prefix" id="mango_prefix" 
                           placeholder="e.g., MANGO-2024-">
                </div>

                <div class="form-group">
                    <label for="prompt">Claude AI Prompt (customize if needed):</label>
                    <textarea name="prompt" id="prompt" rows="10" cols="60">''' + SOPDocumentConfig.DEFAULT_SOP_PROMPT + '''</textarea>
                </div>

                <div class="form-group">
                    <label for="template">Template Filename:</label>
                    <input type="text" name="template" value="LS-sop.docx" required>
                </div>

                <button type="submit">Generate SOP Documents</button>
            </form>

            <div class="status" id="status" style="display: none;"></div>

            <div id="downloads" style="margin-top: 20px;"></div>
        </div>

        <script>
            document.getElementById('uploadForm').addEventListener('submit', async function(e) {
                e.preventDefault();

                const statusDiv = document.getElementById('status');
                const downloadsDiv = document.getElementById('downloads');
                statusDiv.style.display = 'block';
                statusDiv.className = 'status';
                statusDiv.innerHTML = 'Processing files... Please wait.';
                downloadsDiv.innerHTML = '';

                const formData = new FormData(this);

                try {
                    const response = await fetch('/upload', {
                        method: 'POST',
                        body: formData
                    });

                    const result = await response.json();

                    if (result.success) {
                        statusDiv.className = 'status success';
                        statusDiv.innerHTML = `Successfully processed ${result.processed} file(s)!`;

                        // Show download links
                        if (result.output_files && result.output_files.length > 0) {
                            let downloadHTML = '<h3>Download Generated SOPs:</h3><ul>';
                            result.output_files.forEach(file => {
                                downloadHTML += `<li><a href="/download/${file}" download>${file}</a></li>`;
                            });
                            downloadHTML += '</ul>';
                            downloadsDiv.innerHTML = downloadHTML;
                        }
                    } else {
                        throw new Error(result.error || 'Processing failed');
                    }
                } catch (error) {
                    statusDiv.className = 'status error';
                    statusDiv.innerHTML = `Error: ${error.message}`;
                }
            });
        </script>
    </body>
    </html>
    '''


@app.route('/upload', methods=['POST'])
async def upload_files():
    """Handle file upload and processing"""
    files = request.files.getlist('files')
    custom_prompt = request.form.get('prompt', SOPDocumentConfig.DEFAULT_SOP_PROMPT)
    template_name = request.form.get('template', 'LS-sop.docx')
    user_name = request.form.get('user_name', Config.DEFAULT_USER_NAME)
    mango_prefix = request.form.get('mango_prefix', '')

    if not files:
        return jsonify({'error': 'No files uploaded'}), 400

    # Save uploaded files
    uploaded_paths = []
    for file in files:
        if file and file.filename.endswith('.pptx'):
            filename = secure_filename(file.filename)
            filepath = Config.UPLOAD_FOLDER / filename
            file.save(filepath)
            uploaded_paths.append(filepath)

    if not uploaded_paths:
        return jsonify({'error': 'No valid PowerPoint files uploaded'}), 400

    # Process files
    try:
        template_path = Config.TEMPLATE_FOLDER / template_name

        # Create processor with custom configuration
        processor = DocumentProcessor(
            Config.CLAUDE_API_KEY,
            template_path,
            custom_prompt
        )

        # Add user context to processor
        processor.user_name = user_name
        processor.mango_prefix = mango_prefix

        results = await processor.process_batch(uploaded_paths)

        # Clean up uploaded files
        for path in uploaded_paths:
            path.unlink()

        # Return results
        output_files = [str(path.name) for path in results if path]
        return jsonify({
            'success': True,
            'processed': len(results),
            'output_files': output_files
        })

    except Exception as e:
        logger.error(f"Processing error: {str(e)}")
        return jsonify({'error': str(e)}), 500


@app.route('/download/<filename>')
def download_file(filename):
    """Download generated Word document"""
    filepath = Config.OUTPUT_FOLDER / secure_filename(filename)
    if filepath.exists():
        return send_file(filepath, as_attachment=True)
    else:
        return jsonify({'error': 'File not found'}), 404


# API endpoint for Power Apps integration
@app.route('/api/convert', methods=['POST'])
async def api_convert():
    """REST API endpoint for Power Apps integration"""
    data = request.get_json()

    # Validate input
    required_fields = ['file_content', 'filename', 'prompt', 'template']
    for field in required_fields:
        if field not in data:
            return jsonify({'error': f'Missing required field: {field}'}), 400

    # Process the request
    # Implementation depends on how Power Apps sends the file content
    # This is a placeholder for the API logic

    return jsonify({
        'success': True,
        'message': 'API endpoint ready for Power Apps integration'
    })


# Command-line interface
async def main():
    """Command-line interface for batch processing"""
    import argparse

    parser = argparse.ArgumentParser(description='Convert PowerPoint files to Word documents')
    parser.add_argument('--input', '-i', type=str, help='Input directory or file')
    parser.add_argument('--template', '-t', type=str, help='Word template file')
    parser.add_argument('--prompt', '-p', type=str, help='Claude prompt file or text')
    parser.add_argument('--api-key', '-k', type=str, help='Claude API key (or set CLAUDE_API_KEY env var)')
    parser.add_argument('--web', action='store_true', help='Run web interface')

    args = parser.parse_args()

    if args.web:
        # Run Flask app
        app.run(debug=True, host='0.0.0.0', port=5000)
    else:
        # Batch processing mode - now require template and prompt
        if not args.input:
            parser.error("--input is required for batch processing")
        if not args.template:
            parser.error("--template is required for batch processing")
        if not args.prompt:
            parser.error("--prompt is required for batch processing")

        # Set API key
        api_key = args.api_key or Config.CLAUDE_API_KEY
        if not api_key:
            parser.error("Claude API key is required")

        # Load prompt
        prompt_path = Path(args.prompt)
        if prompt_path.exists():
            custom_prompt = prompt_path.read_text()
        else:
            custom_prompt = args.prompt

        # Find PowerPoint files
        input_path = Path(args.input)
        if input_path.is_file():
            pptx_files = [input_path]
        else:
            pptx_files = list(input_path.glob('*.pptx'))

        if not pptx_files:
            print("No PowerPoint files found")
            return

        # Process files
        template_path = Path(args.template)
        processor = DocumentProcessor(api_key, template_path, custom_prompt)

        print(f"Processing {len(pptx_files)} files...")
        results = await processor.process_batch(pptx_files)

        # Summary
        successful = [r for r in results if r]
        print(f"\nCompleted: {len(successful)}/{len(pptx_files)} files processed successfully")
        if successful:
            print("\nGenerated files:")
            for path in successful:
                print(f"  - {path}")


if __name__ == '__main__':
    # Check if running on Windows or WSL
    if sys.platform == 'win32':
        # Windows: Use asyncio.set_event_loop_policy for Windows compatibility
        asyncio.set_event_loop_policy(asyncio.WindowsSelectorEventLoopPolicy())

    # Run the main function
    asyncio.run(main())